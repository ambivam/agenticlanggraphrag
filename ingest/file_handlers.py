from typing import List
import pandas as pd
from docx import Document
from pptx import Presentation
import json
import os
from langchain.docstore.document import Document as LangchainDocument

def process_excel(file_path: str) -> List[LangchainDocument]:
    """Process Excel files (.xlsx, .xls)."""
    try:
        print(f"Processing Excel file: {file_path}")
        # Read all sheets
        df_dict = pd.read_excel(file_path, sheet_name=None)
        documents = []
        
        for sheet_name, df in df_dict.items():
            # Convert dataframe to string representation
            content = f"Sheet: {sheet_name}\n" + df.to_string()
            documents.append(
                LangchainDocument(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "type": "excel",
                        "sheet": sheet_name
                    }
                )
            )
        return documents
    except Exception as e:
        print(f"Error processing Excel file {os.path.basename(file_path)}: {str(e)}")
        import traceback
        print(traceback.format_exc())
        return []

def process_csv(file_path: str) -> List[LangchainDocument]:
    """Process CSV files."""
    try:
        df = pd.read_csv(file_path)
        content = df.to_string()
        return [
            LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "csv"}
            )
        ]
    except Exception as e:
        print(f"Error processing CSV file: {str(e)}")
        return []

def process_word(file_path: str) -> List[LangchainDocument]:
    """Process Word documents (.docx)."""
    try:
        doc = Document(file_path)
        content = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
        
        # Process tables
        for table in doc.tables:
            table_content = []
            for row in table.rows:
                row_content = [cell.text for cell in row.cells]
                table_content.append(" | ".join(row_content))
            content.append("\n".join(table_content))
        
        return [
            LangchainDocument(
                page_content="\n\n".join(content),
                metadata={"source": file_path, "type": "word"}
            )
        ]
    except Exception as e:
        print(f"Error processing Word file: {str(e)}")
        return []

def process_powerpoint(file_path: str) -> List[LangchainDocument]:
    """Process PowerPoint files (.pptx)."""
    try:
        print(f"Processing PowerPoint file: {file_path}")
        prs = Presentation(file_path)
        documents = []
        
        for i, slide in enumerate(prs.slides, 1):
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
            
            if content:
                documents.append(
                    LangchainDocument(
                        page_content="\n".join(content),
                        metadata={
                            "source": file_path,
                            "type": "powerpoint",
                            "slide": i
                        }
                    )
                )
        return documents
    except Exception as e:
        print(f"Error processing PowerPoint file {os.path.basename(file_path)}: {str(e)}")
        import traceback
        print(traceback.format_exc())
        return []

def process_markdown(file_path: str) -> List[LangchainDocument]:
    """Process Markdown files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return [
            LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "markdown"}
            )
        ]
    except Exception as e:
        print(f"Error processing Markdown file: {str(e)}")
        return []

def process_json(file_path: str) -> List[LangchainDocument]:
    """Process JSON files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Convert JSON to formatted string
        content = json.dumps(data, indent=2)
        return [
            LangchainDocument(
                page_content=content,
                metadata={"source": file_path, "type": "json"}
            )
        ]
    except Exception as e:
        print(f"Error processing JSON file: {str(e)}")
        return []

def get_document_loader(file_path: str):
    """Get appropriate document loader based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    
    loaders = {
        '.xlsx': process_excel,
        '.xls': process_excel,
        '.csv': process_csv,
        '.docx': process_word,
        '.pptx': process_powerpoint,
        '.md': process_markdown,
        '.json': process_json
    }
    
    return loaders.get(ext)
